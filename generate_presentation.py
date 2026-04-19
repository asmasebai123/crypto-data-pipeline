#!/usr/bin/env python3
"""
Générateur de présentation pour Crypto Data Pipeline
Crée une présentation professionnelle avec structure complète
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor

# Palettes de couleurs (Ocean Gradient)
COLORS = {
    'deep_blue': RGBColor(6, 90, 130),      # #065A82
    'teal': RGBColor(28, 114, 147),         # #1C7293
    'midnight': RGBColor(33, 41, 60),       # #21295C
    'accent': RGBColor(0, 212, 255),        # #00D4FF
    'white': RGBColor(255, 255, 255),
    'text': RGBColor(51, 51, 51),
    'light_bg': RGBColor(245, 245, 245),
}

def add_title_slide(prs, title, subtitle):
    """Ajoute une slide de titre"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = COLORS['deep_blue']

    # Titre principal
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(9), Inches(1.5))
    title_frame = title_box.text_frame
    title_frame.word_wrap = True
    title_p = title_frame.paragraphs[0]
    title_p.text = title
    title_p.font.size = Pt(54)
    title_p.font.bold = True
    title_p.font.color.rgb = COLORS['white']
    title_p.alignment = PP_ALIGN.CENTER

    # Sous-titre
    subtitle_box = slide.shapes.add_textbox(Inches(0.5), Inches(4.2), Inches(9), Inches(1))
    subtitle_frame = subtitle_box.text_frame
    subtitle_p = subtitle_frame.paragraphs[0]
    subtitle_p.text = subtitle
    subtitle_p.font.size = Pt(24)
    subtitle_p.font.color.rgb = COLORS['accent']
    subtitle_p.alignment = PP_ALIGN.CENTER

    # Auteur
    author_box = slide.shapes.add_textbox(Inches(0.5), Inches(6.8), Inches(9), Inches(0.5))
    author_frame = author_box.text_frame
    author_p = author_frame.paragraphs[0]
    author_p.text = "Asma Sebai | ING2 - Semestre 2 | Avril 2026"
    author_p.font.size = Pt(12)
    author_p.font.italic = True
    author_p.font.color.rgb = COLORS['light_bg']
    author_p.alignment = PP_ALIGN.CENTER

    return slide

def add_content_slide(prs, title):
    """Ajoute une slide de contenu avec titre barré"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = COLORS['white']

    # Barre de titre colorée
    title_bar = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(10), Inches(0.8))
    title_bar.fill.solid()
    title_bar.fill.fore_color.rgb = COLORS['deep_blue']
    title_bar.line.color.rgb = COLORS['deep_blue']

    # Titre
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.15), Inches(9), Inches(0.5))
    title_frame = title_box.text_frame
    title_p = title_frame.paragraphs[0]
    title_p.text = title
    title_p.font.size = Pt(36)
    title_p.font.bold = True
    title_p.font.color.rgb = COLORS['white']
    title_frame.vertical_anchor = MSO_ANCHOR.MIDDLE

    return slide

def add_text(slide, text, left, top, width, height, size=14, bold=False, color=None, align=PP_ALIGN.LEFT):
    """Ajoute du texte à une slide"""
    if color is None:
        color = COLORS['text']

    text_box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    text_frame = text_box.text_frame
    text_frame.word_wrap = True
    text_p = text_frame.paragraphs[0]
    text_p.text = text
    text_p.font.size = Pt(size)
    text_p.font.bold = bold
    text_p.font.color.rgb = color
    text_p.alignment = align
    return text_box

def add_colored_box(slide, text, left, top, width, height, bg_color, text_color, size=12, bold=True):
    """Ajoute une boîte colorée avec texte"""
    shape = slide.shapes.add_shape(1, Inches(left), Inches(top), Inches(width), Inches(height))
    shape.fill.solid()
    shape.fill.fore_color.rgb = bg_color
    shape.line.color.rgb = bg_color

    text_box = shape.text_frame
    text_box.word_wrap = True
    text_p = text_box.paragraphs[0]
    text_p.text = text
    text_p.font.size = Pt(size)
    text_p.font.bold = bold
    text_p.font.color.rgb = text_color
    text_p.alignment = PP_ALIGN.CENTER
    text_box.vertical_anchor = MSO_ANCHOR.MIDDLE

    return shape

# Créer la présentation
prs = Presentation()
prs.slide_width = Inches(10)
prs.slide_height = Inches(7.5)

# ========== SLIDE 1: PAGE DE TITRE ==========
add_title_slide(prs, "Crypto Data Pipeline", "Architecture Data Engineering End-to-End")

# ========== SLIDE 2: PROBLÉMATIQUE ==========
slide = add_content_slide(prs, "Contexte et Problématique")

add_text(slide, "Le Problème", 0.5, 1.2, 4.5, 0.4, size=20, bold=True, color=COLORS['deep_blue'])

problems = [
    "📊 Marché crypto volatil 24h/24 ($billions/jour)",
    "⚠️  Besoin d'analyse en temps réel",
    "🔄 Intégration sources hétérogènes (API, Kafka)",
    "💾 Stockage et transformation complexe",
    "📈 Visualisation accessible décideurs"
]

y = 1.7
for problem in problems:
    add_text(slide, problem, 0.8, y, 8.7, 0.35, size=13)
    y += 0.45

# Boîte objectif
add_colored_box(slide, "OBJECTIFS\n\n✓ Pipeline ETL\n✓ Dual-Engine Pandas/Spark\n✓ Dashboard production\n✓ Cloud 24/7\n✓ Analyses temps réel",
                5.2, 1.1, 4.3, 3, COLORS['teal'], COLORS['white'], size=12)

# ========== SLIDE 3: ARCHITECTURE GLOBALE ==========
slide = add_content_slide(prs, "Architecture Globale du Pipeline")

layers = [
    ("API", 1.2, COLORS['accent']),
    ("INGESTION", 1.9, COLORS['teal']),
    ("STOCKAGE", 2.6, COLORS['deep_blue']),
    ("TRANSFORMATION", 3.3, COLORS['teal']),
    ("ORCHESTRATION", 4.0, COLORS['accent']),
    ("VISUALISATION", 4.7, COLORS['deep_blue'])
]

for name, y, color in layers:
    shape = slide.shapes.add_shape(1, Inches(0.5), Inches(y), Inches(9), Inches(0.5))
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.color.rgb = color

    add_text(slide, name, 0.7, y + 0.08, 8.6, 0.35, size=14, bold=True, color=COLORS['white'])

# Descriptions à côté
descriptions = [
    ("CoinGecko + Kafka", 1.3),
    ("Batch + Streaming", 2.0),
    ("Neon PostgreSQL Cloud", 2.7),
    ("Pandas + PySpark", 3.4),
    ("Prefect + Railway", 4.1),
    ("Streamlit Cloud", 4.8)
]

for desc, y in descriptions:
    add_text(slide, desc, 0.7, y, 2, 0.3, size=11, color=COLORS['text'])

# ========== SLIDE 4: INGESTION BATCH ==========
slide = add_content_slide(prs, "Étape 1: Ingestion Batch (CoinGecko API)")

add_text(slide, "Source: API REST publique CoinGecko", 0.8, 1.4, 4.2, 0.35, size=13)
add_text(slide, "⏱️ Fréquence: Toutes les 10 minutes", 0.8, 1.8, 4.2, 0.35, size=13)
add_text(slide, "📝 Données: Prix, Volume, Variation 24h, Market Cap", 0.8, 2.2, 4.2, 0.35, size=13)
add_text(slide, "🪙 Cryptos: Bitcoin, Ethereum, Solana, XRP, BNB", 0.8, 2.6, 4.2, 0.35, size=13)

# Code example
code_text = """GET /api/v3/coins/markets

Response:
{
  "id": "bitcoin",
  "current_price": 70970.00,
  "market_cap": 1.2T,
  "volume_24h": 29.6B,
  "price_change_24h": -2.84%
}"""

shape = slide.shapes.add_shape(1, Inches(5.2), Inches(1.2), Inches(4.3), Inches(3.8))
shape.fill.solid()
shape.fill.fore_color.rgb = RGBColor(240, 240, 240)
shape.line.color.rgb = COLORS['deep_blue']

add_text(slide, code_text, 5.4, 1.4, 3.9, 3.4, size=9, color=COLORS['text'])

# ========== SLIDE 5: INGESTION STREAMING ==========
slide = add_content_slide(prs, "Étape 2: Ingestion Streaming (Kafka)")

add_text(slide, "Architecture Kafka", 0.5, 1.2, 9, 0.3, size=16, bold=True, color=COLORS['deep_blue'])

# Kafka boxes
kafka_boxes = [
    ("Producer", 0.8),
    ("Topic\ncrypto_prices", 3.5),
    ("Consumer", 6.2)
]

for label, x in kafka_boxes:
    shape = slide.shapes.add_shape(1, Inches(x), Inches(1.8), Inches(2), Inches(1))
    shape.fill.solid()
    shape.fill.fore_color.rgb = COLORS['teal'] if "Topic" not in label else COLORS['accent']
    shape.line.color.rgb = COLORS['deep_blue']
    shape.line.width = Pt(2)

    text_box = shape.text_frame
    text_box.word_wrap = True
    text_p = text_box.paragraphs[0]
    text_p.text = label
    text_p.font.size = Pt(12)
    text_p.font.bold = True
    text_p.font.color.rgb = COLORS['white']
    text_p.alignment = PP_ALIGN.CENTER
    text_box.vertical_anchor = MSO_ANCHOR.MIDDLE

kafka_details = [
    "✓ Flux temps réel 5 secondes",
    "✓ Données simulées (volumes croissants)",
    "✓ Infrastructure scalable pour production",
    "✓ Permet ajout WebSocket futures"
]

y = 3.2
for detail in kafka_details:
    add_text(slide, detail, 0.8, y, 8.2, 0.35, size=13)
    y += 0.45

# ========== SLIDE 6: STOCKAGE ==========
slide = add_content_slide(prs, "Étape 3: Stockage (PostgreSQL Neon Cloud)")

db_tables = [
    ("raw_crypto_prices", "Données brutes batch"),
    ("stream_crypto_prices", "Données streaming temps réel"),
    ("transform_hourly_avg", "Moyennes horaires"),
    ("transform_daily_ranking", "Classements journaliers"),
    ("spark_volatility", "Analyses volatilité (Spark SQL)")
]

y = 1.4
for name, desc in db_tables:
    add_colored_box(slide, name, 0.8, y, 3.2, 0.35, COLORS['teal'], COLORS['white'], size=11)
    add_text(slide, desc, 4.2, y + 0.02, 5.3, 0.3, size=11)
    y += 0.45

add_colored_box(slide, "✓ Gratuit  ✓ Scalable  ✓ Serverless  ✓ Haute disponibilité  ✓ Sauvegarde auto",
                0.8, 4.8, 8.9, 0.5, RGBColor(232, 244, 248), COLORS['text'], size=11, bold=False)

# ========== SLIDE 7: DUAL-ENGINE ==========
slide = add_content_slide(prs, "Étape 4: Transformation (Dual-Engine)")

add_text(slide, "🐼 PANDAS", 0.8, 1.3, 4.2, 0.5, size=16, bold=True, color=COLORS['white'], align=PP_ALIGN.CENTER)
add_colored_box(slide, "🐼 PANDAS", 0.8, 1.3, 4.2, 0.5, COLORS['teal'], COLORS['white'], size=14, bold=True)

add_colored_box(slide, "⚡ PYSPARK", 5.2, 1.3, 4.2, 0.5, COLORS['accent'], COLORS['white'], size=14, bold=True)

pandas_details = [
    "< 50k lignes",
    "Développement rapide",
    "Efficace RAM"
]

spark_details = [
    "≥ 50k lignes",
    "Scalable distribué",
    "Production-grade"
]

y = 2
for detail in pandas_details:
    add_text(slide, detail, 1, y, 3.8, 0.3, size=12)
    y += 0.4

y = 2
for detail in spark_details:
    add_text(slide, detail, 5.4, y, 3.8, 0.3, size=12)
    y += 0.4

add_text(slide, "5 Transformations Implémentées", 0.5, 4.5, 9, 0.3, size=14, bold=True, color=COLORS['deep_blue'])

transformations = [
    "1️⃣ Nettoyage (suppression nulls, doublons, anomalies)",
    "2️⃣ Moyenne horaire (agrégation par crypto/heure)",
    "3️⃣ Classement journalier (Window Functions Spark)",
    "4️⃣ Détection alertes (seuils ±5%)",
    "5️⃣ Volume journalier (somme par crypto/jour)"
]

y = 5
for tr in transformations:
    add_text(slide, tr, 0.8, y, 8.7, 0.25, size=11)
    y += 0.3

# ========== SLIDE 8: ORCHESTRATION ==========
slide = add_content_slide(prs, "Étape 5: Orchestration (Prefect + Railway)")

add_colored_box(slide, "Prefect Flow", 0.8, 1.3, 4.2, 0.4, COLORS['teal'], COLORS['white'], size=12, bold=True)

flow_steps = [
    "1. Créer tables PostgreSQL",
    "2. Appeler API CoinGecko",
    "3. Sauvegarder brutes",
    "4. Charger + Nettoyer",
    "5. Transformation (Pandas/Spark)",
    "6. Analyses SQL avancées",
    "7. Rafraîchir dashboard"
]

y = 1.85
for step in flow_steps:
    add_text(slide, step, 1, y, 3.6, 0.25, size=10)
    y += 0.35

add_colored_box(slide, "Railway Deployment", 5.2, 1.3, 4.2, 0.4, COLORS['accent'], COLORS['white'], size=12, bold=True)

railway_points = [
    "✓ Pipeline tourne 24/7",
    "✓ Exécution toutes les 10min",
    "✓ Gratuit ($5/mois crédit)",
    "✓ Logs en temps réel",
    "✓ Auto-restart en erreur",
    "✓ Scalable production",
    "✓ Zéro maintenance"
]

y = 1.85
for point in railway_points:
    add_text(slide, point, 5.4, y, 3.8, 0.25, size=10)
    y += 0.35

# ========== SLIDE 9: DASHBOARD ==========
slide = add_content_slide(prs, "Étape 6: Visualisation (Streamlit Cloud)")

add_text(slide, "5 Onglets Interactifs", 0.5, 1.2, 9, 0.3, size=16, bold=True, color=COLORS['deep_blue'])

tabs = [
    ("📊 Vue d'ensemble", "KPI cards + Graphiques prix + Classement"),
    ("📈 Analyse technique", "Chandelier OHLC + Heatmap corrélation"),
    ("⚠️  Alertes & Volume", "Volatilité + Volume journalier"),
    ("⚡ Streaming Kafka", "Flux temps réel + Distribution"),
    ("🔧 Spark Analytics", "Volatilité + Dominance marché")
]

y = 1.7
for name, desc in tabs:
    add_colored_box(slide, name, 0.8, y, 2.5, 0.6, COLORS['teal'], COLORS['white'], size=12)
    add_text(slide, desc, 3.6, y + 0.1, 5.2, 0.4, size=11)
    y += 0.75

add_colored_box(slide, "🌐 Streamlit Cloud | 24/7 | Accessible | Auto-refresh 60s", 0.8, 5.8, 8.9, 0.8, COLORS['accent'], COLORS['white'], size=12, bold=True)

# ========== SLIDE 10: RÉSULTATS ==========
slide = add_content_slide(prs, "Résultats Attendus vs Réalité")

add_colored_box(slide, "📊 ATTENDUS", 0.7, 1.35, 4.1, 0.4, RGBColor(232, 244, 248), COLORS['deep_blue'], size=13, bold=True)

expected = [
    "✓ Pipeline ETL",
    "✓ ~5,000 lignes/jour",
    "✓ Dual-Engine",
    "✓ 5 transformations",
    "✓ Dashboard interactive",
    "✓ Orchestration auto",
    "✓ Cloud deployment"
]

y = 1.85
for exp in expected:
    add_text(slide, exp, 0.9, y, 3.9, 0.28, size=11)
    y += 0.35

add_colored_box(slide, "✅ RÉALITÉ", 5.4, 1.35, 3.9, 0.4, RGBColor(232, 248, 232), COLORS['deep_blue'], size=13, bold=True)

reality = [
    "✅ Pipeline complet",
    "✅ 5,050 lignes stockées",
    "✅ Dual-Engine opérationnel",
    "✅ 5 transformations + 3 analyses SQL",
    "✅ Dashboard 5 onglets (v2.1)",
    "✅ Prefect + Railway 24/7",
    "✅ Neon + Streamlit Cloud"
]

y = 1.85
for r in reality:
    add_text(slide, r, 5.4, y, 3.9, 0.28, size=11)
    y += 0.35

# ========== SLIDE 11: MÉTRIQUES ==========
slide = add_content_slide(prs, "Données de Sortie & Métriques")

metrics = [
    ("🪙 Cryptos", "5"),
    ("📊 Lignes/jour", "~5,000"),
    ("⏱️ Batch", "10 min"),
    ("⚡ Stream", "5 sec"),
    ("🗄️ Tables", "8"),
    ("🔄 Transform", "5"),
    ("🔧 Analyses", "3"),
    ("✅ Uptime", "24/7")
]

col = 0
for idx, (label, value) in enumerate(metrics):
    x = 0.8 if col == 0 else 5.2
    y = 1.3 + (idx % 4) * 0.9

    add_colored_box(slide, label + "\n" + value, x, y, 4.2, 0.75, COLORS['teal'], COLORS['white'], size=13)

    if (idx + 1) % 4 == 0:
        col = 1 - col

# ========== SLIDE 12: CONCLUSION ==========
add_title_slide(prs, "Conclusion", "Pipeline Data Engineering Production-Ready")

# Sauvegarder
prs.save("presentation.pptx")
print("✅ Présentation créée: presentation.pptx")
print("📊 14 slides professionnelles générées!")
